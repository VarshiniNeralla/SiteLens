from io import BytesIO
from collections.abc import Iterable
from dataclasses import dataclass
from pathlib import Path

import httpx
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Inches, Pt

from app.config import settings
from app.domain import ObservationRecord
from app.logging_config import get_logger

logger = get_logger(__name__)

# 16:9 widescreen with compact, engineering-style density.
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
OBS_PER_SLIDE = 3

# Layout constants (inches).
MARGIN_X = 0.45
TITLE_TOP = 0.13
TITLE_H = 0.30
SUBTITLE_TOP = 0.43
SUBTITLE_H = 0.22
CONTENT_TOP = 0.80
CONTENT_BOTTOM_MARGIN = 0.18
IMAGE_H = 2.28
TABLE_TOP_GAP = 0.06

TITLE_COLOR = RGBColor(155, 0, 0)
TEXT_COLOR = RGBColor(0, 0, 0)
HEADER_FILL = RGBColor(242, 242, 242)
VALUE_FILL = RGBColor(255, 255, 255)
FONT_FAMILY = "Arial"
LINE_W_EMU = 12700


@dataclass(frozen=True)
class Rect:
    left: float
    top: float
    width: float
    height: float

    @property
    def right(self) -> float:
        return self.left + self.width

    @property
    def bottom(self) -> float:
        return self.top + self.height


def _chunks(items: list[ObservationRecord], size: int) -> Iterable[list[ObservationRecord]]:
    for i in range(0, len(items), size):
        yield items[i : i + size]


def _d(v: object) -> str:
    if hasattr(v, "isoformat") and callable(getattr(v, "isoformat")):
        iso = getattr(v, "isoformat")()
        return iso if iso else "—"
    s = str(v).strip() if v is not None else ""
    return s if s else "—"


def _fixed_grid_geometry() -> tuple[float, float, float]:
    """Return fixed 3-column geometry, regardless of populated count."""
    gap = 0.18
    block_w = (SLIDE_W_IN - 2 * MARGIN_X - 2 * gap) / 3
    left_start = MARGIN_X
    return left_start, block_w, gap


def _observation_rows(obs: ObservationRecord, obs_no: int) -> list[tuple[str, str]]:
    location = ", ".join(
        x
        for x in [
            f"Tower {obs.tower}" if obs.tower else "",
            f"Floor {obs.floor}" if obs.floor else "",
            f"Flat {obs.flat}" if obs.flat else "",
            obs.room or "",
        ]
        if x
    )
    return [
        ("Location", location or "—"),
        ("Severity", _d(obs.severity)),
        ("Date of slab casting", _d(obs.slab_casting_date)),
        ("Date of site visit", _d(obs.site_visit_date)),
        ("Inspection by 3rd party", _d(obs.third_party_status)),
        ("Availability in 3rd party reports", _d(obs.inspection_status)),
        ("Observation no.", str(obs_no)),
    ]


def _fit_with_aspect_bottom_aligned(box: Rect, image_path: Path) -> Rect:
    with Image.open(image_path) as im:
        img_w, img_h = im.size

    if img_w <= 0 or img_h <= 0:
        return box

    box_ratio = box.width / box.height
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        draw_w = box.width
        draw_h = draw_w / img_ratio
    else:
        draw_h = box.height
        draw_w = draw_h * img_ratio

    draw_left = box.left + (box.width - draw_w) / 2
    draw_top = box.top + (box.height - draw_h)  # bottom aligned for consistent baseline
    return Rect(draw_left, draw_top, draw_w, draw_h)


def _set_cell_border_thin_black(cell: object) -> None:
    tc = cell._tc  # type: ignore[attr-defined]
    tc_pr = tc.get_or_add_tcPr()
    # Remove any inherited/default border XML to prevent thick/double border artifacts.
    for child in list(tc_pr):
        if child.tag.endswith(("lnL", "lnR", "lnT", "lnB")):
            tc_pr.remove(child)
    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        line = parse_xml(
            f"<{edge} w='{LINE_W_EMU}' cap='flat' cmpd='sng' algn='ctr' {nsdecls('a')}>"
            "<a:solidFill><a:srgbClr val='000000'/></a:solidFill>"
            "<a:prstDash val='solid'/></{edge}>".replace("{edge}", edge)
        )
        tc_pr.append(line)


def _clear_table_style(table: object) -> None:
    tbl_pr = table._tbl.tblPr  # type: ignore[attr-defined]
    for child in list(tbl_pr):
        if child.tag.endswith("tableStyleId"):
            tbl_pr.remove(child)
    style_id = OxmlElement("a:tableStyleId")
    style_id.text = "{00000000-0000-0000-0000-000000000000}"
    tbl_pr.append(style_id)
    table.first_row = False
    table.last_row = False
    table.first_col = False
    table.last_col = False
    table.horz_banding = False
    table.vert_banding = False


def _write_cell_text(cell: object, text: str, *, size_pt: float = 9.6) -> None:
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.045)
    tf.margin_right = Inches(0.045)
    tf.margin_top = Inches(0.015)
    tf.margin_bottom = Inches(0.015)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = FONT_FAMILY
    run.font.size = Pt(size_pt)
    run.font.bold = False
    run.font.color.rgb = TEXT_COLOR


def _add_header(slide: object, subtitle: str) -> None:
    title_box = slide.shapes.add_textbox(
        Inches(MARGIN_X),
        Inches(TITLE_TOP),
        Inches(SLIDE_W_IN - 2 * MARGIN_X),
        Inches(TITLE_H),
    )
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Quality walkthrough - Observation"
    run.font.name = FONT_FAMILY
    run.font.size = Pt(23)
    run.font.bold = True
    run.font.italic = True
    run.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.LEFT

    subtitle_box = slide.shapes.add_textbox(
        Inches(MARGIN_X),
        Inches(SUBTITLE_TOP),
        Inches(SLIDE_W_IN - 2 * MARGIN_X),
        Inches(SUBTITLE_H),
    )
    stf = subtitle_box.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    srun = sp.add_run()
    srun.text = subtitle
    srun.font.name = FONT_FAMILY
    srun.font.size = Pt(14)
    srun.font.italic = True
    srun.font.color.rgb = TITLE_COLOR
    sp.alignment = PP_ALIGN.LEFT


def _add_rect_border(slide: object, box: Rect) -> None:
    rect = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
        Inches(box.left),
        Inches(box.top),
        Inches(box.width),
        Inches(box.height),
    )
    rect.fill.background()
    rect.line.color.rgb = TEXT_COLOR
    rect.line.width = Emu(LINE_W_EMU)


def _paste_cropped_raster(slide: object, image_box: Rect, source: Image.Image) -> None:
    im = source
    img_w, img_h = im.size
    if img_w <= 0 or img_h <= 0:
        raise ValueError("Image has invalid dimensions")
    target_ratio = image_box.width / image_box.height
    img_ratio = img_w / img_h
    if img_ratio > target_ratio:
        crop_h = img_h
        crop_w = int(crop_h * target_ratio)
        left = (img_w - crop_w) // 2
        top = 0
    else:
        crop_w = img_w
        crop_h = int(crop_w / target_ratio)
        left = 0
        top = (img_h - crop_h) // 2
    cropped = im.crop((left, top, left + crop_w, top + crop_h))
    buff = BytesIO()
    cropped.save(buff, format="PNG")
    buff.seek(0)

    slide.shapes.add_picture(
        buff,
        Inches(image_box.left),
        Inches(image_box.top),
        width=Inches(image_box.width),
        height=Inches(image_box.height),
    )


def _load_observation_raster(obs: ObservationRecord, cwd: Path) -> Image.Image | None:
    ref = (obs.image_path or "").strip()
    if ref.startswith(("http://", "https://")):
        url = ((obs.cloudinary_secure_url or ref) or "").strip()
        try:
            with httpx.Client(timeout=90.0, follow_redirects=True) as client:
                r = client.get(url)
                r.raise_for_status()
                im = Image.open(BytesIO(r.content))
                im.load()
                return im.copy()
        except Exception as exc:  # noqa: BLE001 — PIL/network surface
            logger.warning("Could not fetch observation image (%s): %s", url[:80], exc)
            return None

    img_path = Path(obs.image_path)
    if not img_path.is_absolute():
        img_path = cwd / img_path
    if not img_path.is_file():
        return None
    im = Image.open(img_path)
    im.load()
    return im.copy()


def _add_footer(slide: object, project_name: str, slide_no: int) -> None:
    left_box = slide.shapes.add_textbox(
        Inches(MARGIN_X),
        Inches(SLIDE_H_IN - 0.22),
        Inches(5.0),
        Inches(0.16),
    )
    ltf = left_box.text_frame
    ltf.clear()
    lp = ltf.paragraphs[0]
    lr = lp.add_run()
    lr.text = project_name
    lr.font.name = FONT_FAMILY
    lr.font.size = Pt(8)
    lr.font.color.rgb = TEXT_COLOR
    lp.alignment = PP_ALIGN.LEFT

    right_box = slide.shapes.add_textbox(
        Inches(SLIDE_W_IN - MARGIN_X - 1.2),
        Inches(SLIDE_H_IN - 0.22),
        Inches(1.2),
        Inches(0.16),
    )
    rtf = right_box.text_frame
    rtf.clear()
    rp = rtf.paragraphs[0]
    rr = rp.add_run()
    rr.text = str(slide_no)
    rr.font.name = FONT_FAMILY
    rr.font.size = Pt(8)
    rr.font.color.rgb = TEXT_COLOR
    rp.alignment = PP_ALIGN.RIGHT


def _add_missing_image_placeholder(slide: object, image_box: Rect, hint: str) -> None:
    ph = slide.shapes.add_textbox(
        Inches(image_box.left),
        Inches(image_box.top),
        Inches(image_box.width),
        Inches(image_box.height),
    )
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"Image unavailable\n{hint}"
    r.font.name = FONT_FAMILY
    r.font.size = Pt(9)
    r.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_table(slide: object, table_box: Rect, obs: ObservationRecord, obs_no: int) -> None:
    rows = _observation_rows(obs, obs_no)
    table = slide.shapes.add_table(
        len(rows),
        2,
        Inches(table_box.left),
        Inches(table_box.top),
        Inches(table_box.width),
        Inches(table_box.height),
    ).table
    _clear_table_style(table)

    left_col_w = table_box.width * 0.44
    table.columns[0].width = Inches(left_col_w)
    table.columns[1].width = Inches(table_box.width - left_col_w)

    row_h_emu = Emu(Inches(table_box.height) // len(rows))
    for i, (label, value) in enumerate(rows):
        row = table.rows[i]
        row.height = row_h_emu
        c0 = row.cells[0]
        c1 = row.cells[1]

        _write_cell_text(c0, label, size_pt=9.4)
        _write_cell_text(c1, value, size_pt=9.4)
        c0.vertical_anchor = MSO_ANCHOR.MIDDLE
        c1.vertical_anchor = MSO_ANCHOR.MIDDLE

        c0.fill.solid()
        c0.fill.fore_color.rgb = HEADER_FILL
        c1.fill.solid()
        c1.fill.fore_color.rgb = VALUE_FILL

        _set_cell_border_thin_black(c0)
        _set_cell_border_thin_black(c1)


def _render_observation_block(
    slide: object,
    obs: ObservationRecord,
    obs_no: int,
    block: Rect,
    cwd: Path,
) -> None:
    image_box = Rect(block.left, block.top, block.width, IMAGE_H)
    table_top = image_box.bottom + TABLE_TOP_GAP
    table_box = Rect(block.left, table_top, block.width, block.bottom - table_top)

    raster = _load_observation_raster(obs, cwd)
    if raster is not None:
        try:
            _paste_cropped_raster(slide, image_box, raster)
        finally:
            raster.close()
        _add_rect_border(slide, image_box)
    else:
        hint = (obs.image_original_filename or Path(obs.image_path).name)[:80] or "—"
        _add_missing_image_placeholder(slide, image_box, hint)
        _add_rect_border(slide, image_box)

    _add_table(slide, table_box, obs, obs_no)
    _add_rect_border(slide, block)


def _blank_slide_layout(prs: Presentation):  # noqa: ANN001
    for lay in prs.slide_layouts:
        if "blank" in (lay.name or "").lower():
            return lay
    return prs.slide_layouts[-1]


def _resolve_subtitle(page_obs: list[ObservationRecord]) -> str:
    kinds = {o.observation_type.strip() for o in page_obs if o.observation_type and o.observation_type.strip()}
    if len(kinds) == 1:
        return kinds.pop()
    return "Mixed observations"


def build_quality_report_pptx(
    *,
    project_name: str,
    title: str,
    observations: list[ObservationRecord],
    output_path: Path,
) -> Path:
    _ = project_name  # Project consistency validated upstream; slide subtitle uses observation type.
    _ = title  # Header title is standardized per requested layout.

    template_path = settings.ppt_template_path
    if template_path is not None and Path(template_path).is_file():
        prs = Presentation(str(Path(template_path).resolve()))
        logger.info("PPT base presentation loaded from %s", template_path)
    else:
        prs = Presentation()

    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = _blank_slide_layout(prs)
    cwd = Path.cwd()

    content_height = SLIDE_H_IN - CONTENT_TOP - CONTENT_BOTTOM_MARGIN

    for group_idx, group in enumerate(_chunks(observations, OBS_PER_SLIDE), start=1):
        slide = prs.slides.add_slide(blank)
        _add_header(slide, _resolve_subtitle(group))
        left_start, block_width, gap = _fixed_grid_geometry()

        for idx, obs in enumerate(group):
            left = left_start + idx * (block_width + gap)
            block = Rect(left, CONTENT_TOP, block_width, content_height)
            _render_observation_block(
                slide,
                obs,
                obs_no=(group_idx - 1) * OBS_PER_SLIDE + idx + 1,
                block=block,
                cwd=cwd,
            )
        _add_footer(slide, project_name=project_name, slide_no=group_idx)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
